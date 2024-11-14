import sys
from PyQt5.QtWidgets import QApplication, QMainWindow, QVBoxLayout, QWidget, QTextEdit, QPushButton, QToolBar, QLabel, QHBoxLayout, QFileDialog
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QFont
from groq import Groq
from fpdf import FPDF
import docx
from PyPDF2 import PdfReader
from pptx import Presentation

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()

        self.setWindowTitle("Simplification of Concepts")
        self.setGeometry(100, 100, 800, 600)

        # Main layout
        main_layout = QVBoxLayout()

        # Toolbar
        toolbar = QToolBar()
        toolbar.setStyleSheet("background-color: black; height: 10%;")
        self.addToolBar(toolbar)

        # LOGO
        logo = QLabel("BreakDown")
        logo.setStyleSheet("color: white; font-size: 24px; font-weight: bold;")
        toolbar.addWidget(logo)

        # File open button
        self.open_file_button = QPushButton("Open File")
        self.open_file_button.setStyleSheet("background-color: black; color: white; font-size: 16px; padding: 10px;")
        self.open_file_button.clicked.connect(self.open_file)
        toolbar.addWidget(self.open_file_button)

        # Submit button
        self.submit_button = QPushButton("Submit")
        self.submit_button.setStyleSheet("background-color: black; color: white; font-size: 16px; padding: 10px;")
        self.submit_button.clicked.connect(self.call_api)
        toolbar.addWidget(self.submit_button)

        #Download button
        self.download_button = QPushButton("Download")
        self.download_button.setStyleSheet("background-color: black; color: white; font-size: 16px; padding: 10px;")
        self.download_button.clicked.connect(self.download_pdf)
        toolbar.addWidget(self.download_button)

        # Central widget
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        central_widget.setLayout(main_layout)

        # Input layout
        input_layout = QHBoxLayout()

        # Text area for input
        self.input_area = QTextEdit()
        self.input_area.setPlaceholderText("Enter your text here...")
        self.input_area.setStyleSheet("font-size: 16px; padding: 10px;")
        self.input_area.setFixedHeight(200)  # Set initial height to half
        input_layout.addWidget(self.input_area)

        main_layout.addLayout(input_layout)

        # Output field
        self.output_field = QTextEdit()
        self.output_field.setReadOnly(True)
        self.output_field.setStyleSheet("font-size: 16px; padding: 10px;")
        main_layout.addWidget(self.output_field)

        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(20)

    def call_api(self):
        user_input = self.input_area.toPlainText()
        if user_input:
            client = Groq(api_key="")
            chat_completion = client.chat.completions.create(
                messages=[
                    {
                        "role": "user",
                        "content": user_input,
                    }
                ],
                model="llama3-8b-8192",
            )
            response = chat_completion.choices[0].message.content
            self.output_field.setText(response)

    def open_file(self):
        options = QFileDialog.Options()
        file_name, _ = QFileDialog.getOpenFileName(self, "Open File", "", "All Files (*);;Text Files (*.txt);;PDF Files (*.pdf);;Word Files (*.docx);;PowerPoint Files (*.pptx)", options=options)
        if file_name:
            if file_name.endswith('.txt'):
                with open(file_name, 'r') as file:
                    self.input_area.setText(file.read())
            elif file_name.endswith('.pdf'):
                with open(file_name, 'rb') as file:
                    reader = PdfReader(file)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text()
                    self.input_area.setText(text)
            elif file_name.endswith('.docx'):
                doc = docx.Document(file_name)
                text = "\n".join([para.text for para in doc.paragraphs])
                self.input_area.setText(text)
            elif file_name.endswith('.pptx'):
                prs = Presentation(file_name)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                self.input_area.setText(text)

    def download_pdf(self):
        text = self.output_field.toPlainText()
        if text:
            try:
                options = QFileDialog.Options()
                file_name, _ = QFileDialog.getSaveFileName(self, "Save PDF", "", "PDF Files (*.pdf);;All Files (*)", options=options)
                if file_name:
                    pdf = FPDF()
                    pdf.add_page()
                    pdf.set_font("Arial", size=12)
                    pdf.multi_cell(0, 10, text)
                    pdf.output(file_name)
                    print("PDF generated successfully.")
            except Exception as e:
                print(f"Error generating PDF: {e}")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())